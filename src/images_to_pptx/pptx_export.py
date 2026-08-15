from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from images_to_pptx.slides import list_slide_images

EMU_PER_PX = 9525


def export_pptx(directory: Path, output: Path) -> int:
    images = list_slide_images(directory)
    if not images:
        raise FileNotFoundError("Нет изображений slide-N.png в выбранной папке")
    with Image.open(images[0]) as first:
        width_px, height_px = first.size
    prs = Presentation()
    prs.slide_width = Emu(width_px * EMU_PER_PX)
    prs.slide_height = Emu(height_px * EMU_PER_PX)
    blank = prs.slide_layouts[6]
    for image_path in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            Emu(0),
            Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return len(images)
